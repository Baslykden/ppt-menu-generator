import sys
from pathlib import Path
from tkinter import Tk, filedialog, messagebox
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN


# ---- Константы верстки (как в VBA) ----
MAX_SCAN_ROWS = 400
MAX_TABLE_HEIGHT_CM = 14.8       # максимальная высота таблицы на слайде
ROW_HEIGHT_HEADER_CM = 1.92      # заголовок таблицы
ROW_HEIGHT_DATA_CM = 0.7         # строка данных
MAIN_HEIGHT_LAST_CM = 12.1       # лимит высоты данных на последнем слайде, если там ещё итоги

COL1_WIDTH_CM = 15.0
COL2_WIDTH_CM = 2.9
COL3_WIDTH_CM = 2.9
COL4_WIDTH_CM = 3.5


def pick_excel_file() -> Path:
    """Диалог выбора Excel-файла."""
    root = Tk()
    root.withdraw()

    messagebox.showinfo(
        "Генератор презентации",
        "Выберите Excel-файл, из которого нужно сделать презентацию."
    )

    filetypes = [
        ("Excel files", "*.xlsx *.xlsm"),
        ("All files", "*.*"),
    ]
    filename = filedialog.askopenfilename(
        title="Выбор Excel-файла",
        filetypes=filetypes,
    )

    if not filename:
        sys.exit(0)

    return Path(filename)


def load_book(path: Path):
    """Открываем книгу openpyxl с готовыми значениями формул."""
    return load_workbook(path, data_only=True)


def get_skip_columns_flag(wb):
    """Флаг пропуска столбцов (лист 'Расчет стоимости', ячейка I1)."""
    try:
        ws = wb["Расчет стоимости"]
    except KeyError:
        return False

    val = ws["I1"].value
    return bool(val)


def get_headers_sheet(wb):
    """Лист с заголовками и служебными данными (11-й лист в книге)."""
    try:
        return wb.worksheets[10]  # индекс 10 => 11-й лист (как Worksheets(11))
    except IndexError:
        raise RuntimeError("В книге должен быть минимум 11 листов (для заголовков и настроек).")


def get_header_text(wb, ws, sheet_index: int):
    """
    Заголовок слайда:
    C1 текущего листа + G(индекс-2) на листе 1 + A2 на листе 11.
    """
    ws1 = wb.worksheets[0]
    ws11 = get_headers_sheet(wb)

    c1 = ws["C1"].value or ""
    g_val = ws1[f"G{sheet_index - 2}"].value or ""
    a2 = ws11["A2"].value or ""

    # аккуратно собираем, чтобы не было лишних запятых/пробелов
    parts = [str(c1)]
    if g_val:
        parts.append(str(g_val))
    if a2:
        parts.append(str(a2))
    return ", ".join(parts[:-1]) + (" " + parts[-1] if len(parts) > 1 else parts[0])


def collect_rows_for_sheet(wb, ws, skip_columns: bool):
    """
    Собираем данные с листа:
    B – категория, C – название, D – вес порции, E – порции, F – г/чел.
    Режем заголовочные/служебные строки, сортируем по категории.
    Возвращаем список строк: [category, name, weight, portions, g_per_person].
    """
    ws11 = get_headers_sheet(wb)
    hdr_w = str(ws11["A1"].value or "")
    hdr_p = str(ws11["B1"].value or "")
    hdr_g = str(ws11["C1"].value or "")

    rows_raw = []

    last_row = min(ws.max_row, MAX_SCAN_ROWS)

    for i in range(2, last_row + 1):
        cell_E = ws[f"E{i}"].value
        if cell_E is None or cell_E == 0:
            continue

        cat_cell = str(ws[f"B{i}"].value or "")
        name_cell = str(ws[f"C{i}"].value or "")
        d_val = ws[f"D{i}"].value
        e_val = ws[f"E{i}"].value
        f_val = ws[f"F{i}"].value

        # пропускаем заголовочные строки
        if ("Категория блюд" in cat_cell) or ("Наименован" in name_cell) \
                or (str(d_val) == hdr_w) or (str(e_val) == hdr_p) or (str(f_val) == hdr_g):
            continue

        category = cat_cell
        name = name_cell
        weight = d_val if not skip_columns else None
        portions = e_val if not skip_columns else None
        g_per_person = f_val

        rows_raw.append([category, name, weight, portions, g_per_person])

    if not rows_raw:
        return []

    # сортировка по категории (1-й столбец)
    rows_raw.sort(key=lambda r: str(r[0] or ""))

    return rows_raw


def build_master_rows_and_totals(wb, rows_raw, skip_columns: bool):
    """
    Из отсортированных строк строим master_rows:
      (is_category, text, weight, portions, g_per_person)
    плюс считаем totalFoodPerPerson / totalLiquidPerPerson.
    """
    ws11 = get_headers_sheet(wb)

    # список "жидких" категорий A8:A12
    valid_categories = set()
    for row in ws11["A8:A12"]:
        for cell in row:
            if cell.value not in (None, ""):
                valid_categories.add(str(cell.value))

    # --- считаем границы категорий ---
    categories = []
    i = 0
    n = len(rows_raw)
    while i < n:
        cat = str(rows_raw[i][0] or "")
        start = i
        j = i + 1
        while j < n and str(rows_raw[j][0] or "") == cat:
            j += 1
        end = j - 1
        categories.append((cat, start, end))
        i = j

    # --- считаем итоги ---
    total_food = 0.0
    total_liquid = 0.0

    for row in rows_raw:
        cat_name = str(row[0] or "")
        val = row[4]
        if isinstance(val, (int, float)):
            if cat_name in valid_categories:
                total_liquid += float(val)
            else:
                total_food += float(val)

    # --- master_rows ---
    master_rows = []

    for cat, start, end in categories:
        if end < start:
            continue

        # строка категории
        master_rows.append((True, cat, None, None, None))

        # блюда
        for idx in range(start, end + 1):
            _, name, weight, portions, gpp = rows_raw[idx]
            master_rows.append(
                (False, str(name or ""),
                 weight if not skip_columns else None,
                 portions if not skip_columns else None,
                 gpp)
            )

    return master_rows, total_food, total_liquid


def split_master_rows_to_slides(master_rows):
    """
    Делим master_rows на слайды по MAX_ROWS_PER_SLIDE.
    """
    max_rows_per_slide = int(
        (MAX_TABLE_HEIGHT_CM - ROW_HEIGHT_HEADER_CM) / ROW_HEIGHT_DATA_CM
    )  # как в VBA (≈18)

    slides = []
    i = 0
    n = len(master_rows)
    while i < n:
        slides.append(master_rows[i:i + max_rows_per_slide])
        i += max_rows_per_slide
    return slides

def remove_table_borders(table):
    """Полностью убирает линии (границы) у таблицы."""
    # убираем границы у каждой ячейки
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for ln_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr"):
                ln = tcPr.find(qn(ln_tag))
                if ln is not None:
                    tcPr.remove(ln)

    # убираем общие границы таблицы (tblBorders), если есть
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        borders = tblPr.find(qn("a:tblBorders"))
        if borders is not None:
            tblPr.remove(borders)

def apply_no_grid_style(shape):
    """
    Присваивает таблице встроенный стиль PowerPoint:
    'No Style, No Grid' = {2D5ABB26-0587-4C30-8999-92F81FD0307C}
    как в VBA: tbl.ApplyStyle "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    """
    try:
        tbl = shape._element.graphic.graphicData.tbl
    except AttributeError:
        return  # на всякий случай, если shape не таблица

    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = OxmlElement('a:tblPr')
        tbl.append(tblPr)

    guid = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

    # ищем существующий <a:tableStyleId>, если есть — переписываем
    for child in list(tblPr.iterchildren()):
        if child.tag == qn('a:tableStyleId'):
            child.text = guid
            return

    # если не нашли — добавляем новый
    tsid = OxmlElement('a:tableStyleId')
    tsid.text = guid
    tblPr.append(tsid)

def clear_table_borders(table):
    """
    Удаляем все линии границ у таблицы (делаем сетку полностью прозрачной).
    """
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            # Все возможные линии: слева, справа, сверху, снизу и диагонали
            for border_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr"):
                ln = tcPr.find(qn(border_tag))
                if ln is not None:
                    tcPr.remove(ln)

def set_table_style_no_grid(table):
    """
    Задаём для таблицы стиль PowerPoint: 'No Style, No Grid'
    GUID тот же, что ты использовал в VBA: {2D5ABB26-0587-4C30-8999-92F81FD0307C}
    """
    tbl = table._tbl

    tblPr_xml = """
    <a:tblPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
             xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
             xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
             firstRow="1" bandRow="1">
      <a:tableStyleId>{2D5ABB26-0587-4C30-8999-92F81FD0307C}</a:tableStyleId>
    </a:tblPr>
    """
    new_tblPr = parse_xml(tblPr_xml)

    # удаляем старый tblPr (если был)
    existing_tblPr = tbl.xpath('./a:tblPr')
    if existing_tblPr:
        tbl.remove(existing_tblPr[0])

    # вставляем новый tblPr в начало
    tbl.insert(0, new_tblPr)


def create_slide_with_table(
    prs,
    header_text,
    bg_image_path,
    slide_rows,
    skip_columns,
    is_last_slide,
    total_food_per_person,
    total_liquid_per_person,
):
    """Создаём слайд и рисуем на нём таблицу для переданных строк."""
    data_rows = len(slide_rows)
    extra_rows = 3 if is_last_slide else 0      # пустая + 2 итоговые
    total_rows = 1 + data_rows + extra_rows     # +1 строка заголовков

    # --- создаём слайд ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # пустой макет

    # фон
    if bg_image_path.exists():
        slide.shapes.add_picture(
            str(bg_image_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # заголовок события
    tb = slide.shapes.add_textbox(Cm(1.5), Cm(1.0), prs.slide_width - Cm(3), Cm(1.8))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = header_text
    p.font.name = "Century Gothic"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    # --- параметры таблицы ---
    table_total_width = Cm(
        COL1_WIDTH_CM + COL2_WIDTH_CM + COL3_WIDTH_CM + COL4_WIDTH_CM
    )
    left = int((prs.slide_width - table_total_width) / 2)
    top = Cm(3.2)
    height = Cm(MAX_TABLE_HEIGHT_CM)

    shape = slide.shapes.add_table(total_rows, 4, left, top, table_total_width, height)
    table = shape.table

    # ВАЖНО: стиль "No Style, No Grid" — убирает сетку и фон, как в VBA
    set_table_style_no_grid(table)

    # ширина столбцов
    table.columns[0].width = Cm(COL1_WIDTH_CM)
    table.columns[1].width = Cm(COL2_WIDTH_CM)
    table.columns[2].width = Cm(COL3_WIDTH_CM)
    table.columns[3].width = Cm(COL4_WIDTH_CM)

    # высота строк
    table.rows[0].height = Cm(ROW_HEIGHT_HEADER_CM)
    for r in range(1, total_rows):
        table.rows[r].height = Cm(ROW_HEIGHT_DATA_CM)

    # общий шрифт + прозрачная заливка по всем ячейкам
    for r in range(total_rows):
        for c in range(4):
            cell = table.cell(r, c)
            cell.fill.background()  # фон = прозрачный
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT

    # лист с заголовками/подписями
    ws11 = getattr(prs, "_headers_ws", None)

    hdr_w = "Вес порции, грамм"
    hdr_p = "Кол-во порций"
    hdr_g = "Вес на одну персону, грамм"
    label_food = "Итого по еде"
    label_liquid = "Итого по напиткам"

    if ws11 is not None:
        hdr_w = str(ws11["A1"].value or hdr_w)
        hdr_p = str(ws11["B1"].value or hdr_p)
        hdr_g = str(ws11["C1"].value or hdr_g)
        label_food = str(ws11["A4"].value or label_food)
        label_liquid = str(ws11["A5"].value or label_liquid)

    # --- заголовки столбцов ---
    cell = table.cell(0, 0)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Наименования блюд"
    p.font.bold = True
    p.font.name = "Century Gothic"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    if skip_columns:
        table.cell(0, 1).text = ""
        table.cell(0, 2).text = ""
    else:
        cell = table.cell(0, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = hdr_w
        p.font.bold = True
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        cell = table.cell(0, 2)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = hdr_p
        p.font.bold = True
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    cell = table.cell(0, 3)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = hdr_g
    p.font.bold = True
    p.font.name = "Century Gothic"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # --- данные: категории + блюда ---
    for idx, row in enumerate(slide_rows, start=0):
        is_category, text, weight, portions, gpp = row
        row_idx = 1 + idx  # +1 из-за строки заголовков

        # первый столбец
        cell = table.cell(row_idx, 0)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text or ""
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.bold = bool(is_category)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT
        # отступ для наименований
        if is_category:
            tf.margin_left = Cm(0)
        else:
            tf.margin_left = Cm(0.35)

        # остальные столбцы
        if is_category:
            w_text = ""
            q_text = ""
            g_text = ""
        else:
            w_text = "" if (skip_columns or weight is None) else str(weight)
            q_text = "" if (skip_columns or portions is None) else str(portions)
            if isinstance(gpp, (int, float)):
                g_text = f"{float(gpp):.2f}".replace(".", ",")
            else:
                g_text = "" if gpp is None else str(gpp)

        # столбец 2
        cell = table.cell(row_idx, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = w_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        # столбец 3
        cell = table.cell(row_idx, 2)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = q_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        # столбец 4
        cell = table.cell(row_idx, 3)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = g_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.RIGHT

    # --- блок итогов ---
    if is_last_slide:
        total_rows = len(table.rows)
        if total_rows >= 4:
            row_blank = total_rows - 3   # пустая
            row_food = total_rows - 2    # итого еда
            row_liquid = total_rows - 1  # итого напитки

            # очищаем пустую строку
            for c in range(4):
                table.cell(row_blank, c).text = ""

            # ИТОГО по еде
            cell = table.cell(row_food, 0)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label_food + ":"
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            tf.margin_left = Cm(0)

            cell = table.cell(row_food, 3)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{total_food_per_person:.2f}".replace(".", ",")
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.RIGHT

            # ИТОГО по напиткам
            cell = table.cell(row_liquid, 0)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label_liquid + ":"
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            tf.margin_left = Cm(0)

            cell = table.cell(row_liquid, 3)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{total_liquid_per_person:.2f}".replace(".", ",")
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.RIGHT


def process_sheet(
    wb,
    ws,
    sheet_index: int,
    prs: Presentation,
    bg_image_path: Path,
    skip_columns: bool,
):
    """Полный цикл для одного листа: сбор данных, разбиение на слайды, создание слайдов."""

    # 1. Собираем строки для листа
    rows_raw = collect_rows_for_sheet(wb, ws, skip_columns)
    if not rows_raw:
        return

    # 2. Строим общий список строк и считаем итоги по еде/жидкости
    master_rows, total_food_per_person, total_liquid_per_person = build_master_rows_and_totals(
        wb, rows_raw, skip_columns
    )

    if not master_rows:
        return

    # 3. Разбиваем master_rows на логические "страницы" (по количеству строк)
    slides = split_master_rows_to_slides(master_rows)
    if not slides:
        return

    # 4. Проверяем, влезут ли итоговые строки на последний слайд
    last_rows_count = len(slides[-1])
    main_height_last = ROW_HEIGHT_HEADER_CM + last_rows_count * ROW_HEIGHT_DATA_CM
    can_place_totals_on_last = main_height_last <= MAIN_HEIGHT_LAST_CM

    # 5. Заголовок мероприятия для всех слайдов этого листа
    header_text = get_header_text(wb, ws, sheet_index)

    # 6. Рисуем слайды с таблицами
    for idx, slide_rows in enumerate(slides, start=1):
        # "последний" только если это реально последний слайд И итоги туда влезают
        is_last = (idx == len(slides)) and can_place_totals_on_last

        # ВАЖНО: вызываем БЕЗ именованных аргументов, только позиционно
        create_slide_with_table(
            prs,
            header_text,
            bg_image_path,
            slide_rows,
            skip_columns,
            is_last,
            total_food_per_person,
            total_liquid_per_person,
        )

    # 7. Если итоги на последний не влезают — отдельный слайд только под итоги
    if not can_place_totals_on_last:
        create_slide_with_table(
            prs,
            header_text,
            bg_image_path,
            [],                # без строк блюд, только заголовок + итоги
            skip_columns,
            True,              # это слайд с итогами
            total_food_per_person,
            total_liquid_per_person,
        )


def main():
    excel_path = pick_excel_file()
    wb = load_book(excel_path)

    bg_image_path = excel_path.parent / "image.png"
    if not bg_image_path.exists():
        messagebox.showwarning(
            "Внимание",
            f"Файл фона image.png не найден в папке:\n{excel_path.parent}\n"
            f"Слайды будут без фона."
        )

    prs = Presentation()
    # прикрепим лист с заголовками к объекту презентации, чтобы
    # удобно доставать его внутри create_slide_with_table
    prs._headers_ws = get_headers_sheet(wb)  # type: ignore

    skip_columns = get_skip_columns_flag(wb)

    # обработка листов 3..8 (индекс 3–8 в Excel => 2..7 в списке)
    for idx, ws in enumerate(wb.worksheets, start=1):
        if 3 <= idx <= 8:
            process_sheet(wb, ws, idx, prs, bg_image_path, skip_columns)

    out_path = excel_path.with_name(excel_path.stem + "_auto.pptx")
    prs.save(out_path)

    messagebox.showinfo(
        "Готово",
        f"Презентация сохранена:\n{out_path}"
    )


if __name__ == "__main__":
    main()
